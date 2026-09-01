import hashlib
import io
import json
import sys
import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from pptx import Presentation
from sqlalchemy import create_engine
from sqlalchemy.orm import sessionmaker
from starlette.datastructures import UploadFile


BACKEND_DIR = Path(__file__).resolve().parents[1] / "backend"
sys.path.insert(0, str(BACKEND_DIR))

from capabilities.tools import ppt_native  # noqa: E402
import agent  # noqa: E402
import main  # noqa: E402
from database import Base  # noqa: E402
from harness.tools import ToolContext  # noqa: E402
from models import AttachmentModel, ConversationModel, RunModel, UserModel  # noqa: E402


def create_pptx(path: Path) -> None:
    deck = Presentation()
    first = deck.slides.add_slide(deck.slide_layouts[1])
    first.shapes.title.text = "旧标题"
    first.placeholders[1].text = "旧正文"
    second = deck.slides.add_slide(deck.slide_layouts[1])
    second.shapes.title.text = "第二页"
    second.placeholders[1].text = "第二页正文"
    deck.save(path)


class NativePptToolsTests(unittest.IsolatedAsyncioTestCase):
    async def asyncSetUp(self):
        self.engine = create_engine("sqlite:///:memory:")
        Base.metadata.create_all(self.engine)
        self.sessions = sessionmaker(bind=self.engine)
        self.temp = tempfile.TemporaryDirectory()
        self.root = Path(self.temp.name)
        self.source = self.root / "source.pptx"
        create_pptx(self.source)
        self.original_hash = hashlib.sha256(self.source.read_bytes()).hexdigest()
        self.conversation_id = "00000000-0000-0000-0000-000000000501"
        db = self.sessions()
        user = UserModel(username="native-ppt-user")
        db.add(user)
        db.flush()
        db.add(ConversationModel(id=self.conversation_id, user_id=user.id, title="PPTX"))
        attachment = AttachmentModel(
            conversation_id=self.conversation_id,
            file_name="source.pptx",
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            size_bytes=self.source.stat().st_size,
            storage_path=str(self.source),
        )
        db.add(attachment)
        db.commit()
        self.attachment_id = attachment.id
        db.close()
        self.context = ToolContext(conversation_id=self.conversation_id)
        self.patches = (
            patch.object(ppt_native, "SessionLocal", self.sessions),
            patch.object(ppt_native, "NATIVE_STORAGE_ROOT", self.root / "native"),
        )
        for item in self.patches:
            item.start()

    async def asyncTearDown(self):
        for item in reversed(self.patches):
            item.stop()
        self.temp.cleanup()

    async def test_template_fill_is_planned_applied_and_does_not_overwrite_source(self):
        analyzed = await ppt_native.execute_analyze_pptx_template(
            {"source_id": self.attachment_id, "slides": "1"}, self.context
        )
        payload = json.loads(analyzed.content)
        plan = payload["draft_plan"]
        title_slot = next(item for item in plan["slides"][0]["replacements"] if item["old_text"] == "旧标题")
        title_slot["text"] = "新标题"

        prepared = await ppt_native.execute_prepare_pptx_template_fill(
            {"project_id": payload["project_id"], "fill_plan": plan, "accept_warnings": True},
            self.context,
        )
        prepared_payload = json.loads(prepared.content)
        self.assertTrue(prepared_payload["ready_for_approval"])

        applied = await ppt_native.execute_apply_pptx_template_fill(
            {"project_id": payload["project_id"]}, self.context
        )
        artifact = applied.data["artifacts"][0]
        output = Path(artifact["file_path"])
        self.assertTrue(output.exists())
        self.assertEqual(hashlib.sha256(self.source.read_bytes()).hexdigest(), self.original_hash)
        generated = Presentation(output)
        self.assertEqual(generated.slides[0].shapes.title.text, "新标题")

    async def test_native_enhancement_adds_notes_and_transition_without_overwriting_source(self):
        prepared = await ppt_native.execute_prepare_pptx_enhancement(
            {
                "source_id": self.attachment_id,
                "title": "增强版",
                "notes": [
                    {"slide": 1, "text": "第一页讲稿"},
                    {"slide": 2, "text": "第二页讲稿"},
                ],
                "transition": "fade",
                "transition_duration": 0.6,
            },
            self.context,
        )
        payload = json.loads(prepared.content)
        self.assertTrue(payload["ready_for_approval"])

        applied = await ppt_native.execute_apply_pptx_enhancement(
            {"project_id": payload["project_id"]}, self.context
        )
        output = Path(applied.data["artifacts"][0]["file_path"])
        self.assertTrue(output.exists())
        self.assertEqual(hashlib.sha256(self.source.read_bytes()).hexdigest(), self.original_hash)
        with __import__("zipfile").ZipFile(output) as package:
            names = set(package.namelist())
            self.assertIn("ppt/notesSlides/notesSlide1.xml", names)
            slide_xml = package.read("ppt/slides/slide1.xml").decode("utf-8")
            self.assertIn("transition", slide_xml)

    async def test_project_cannot_be_used_from_another_conversation(self):
        analyzed = await ppt_native.execute_analyze_pptx_template(
            {"source_id": self.attachment_id, "slides": "1"}, self.context
        )
        project_id = json.loads(analyzed.content)["project_id"]
        with self.assertRaises(ValueError):
            await ppt_native.execute_apply_pptx_template_fill(
                {"project_id": project_id}, ToolContext(conversation_id="another-conversation")
            )

    async def test_binary_pptx_artifact_is_copied_with_markdown_preview(self):
        controlled = self.root / "storage"
        controlled.mkdir()
        source = controlled / "native-output.pptx"
        create_pptx(source)
        db = self.sessions()
        run = RunModel(conversation_id=self.conversation_id, status="running")
        db.add(run)
        db.commit()
        with (
            patch.object(agent, "STORAGE_DIR", str(controlled)),
            patch.object(agent, "pptx_pdf_preview_available", return_value=False),
            patch("artifact_display.libreoffice_executable", return_value=None),
            patch.dict("os.environ", {
                "LIBREOFFICE_CONVERT_URL": "",
                "LIBREOFFICE_UNOSERVER_HOST": "",
            }),
        ):
            artifact, payload = agent._persist_artifact(db, run.id, {
                "artifact_type": "ppt",
                "title": "原生PPTX",
                "mime_type": ppt_native.PPTX_MIME,
                "extension": "pptx",
                "file_path": str(source),
                "preview_kind": "markdown",
                "preview_content": "# 原生PPTX\n\n已通过校验。",
            }, 0)
            self.assertEqual(Path(artifact.storage_path).read_bytes(), source.read_bytes())
            self.assertEqual(payload["markdown"], "# 原生PPTX\n\n已通过校验。")
            self.assertEqual(agent.read_artifact_preview(artifact)[0], "markdown")
        db.close()

    async def test_upload_accepts_pptx_and_numbered_audio_without_text_parsing(self):
        with (
            patch.object(main, "SessionLocal", self.sessions),
            patch.object(main, "ATTACHMENT_DIR", str(self.root / "attachments")),
        ):
            (self.root / "attachments").mkdir()
            pptx_response = await main.upload_file(
                UploadFile(file=io.BytesIO(self.source.read_bytes()), filename="template.pptx"),
                self.conversation_id,
            )
            audio_response = await main.upload_file(
                UploadFile(file=io.BytesIO(b"audio"), filename="001.mp3"),
                self.conversation_id,
            )
        self.assertTrue(json.loads(pptx_response.body)["success"])
        self.assertTrue(json.loads(audio_response.body)["success"])
        db = self.sessions()
        uploaded = db.query(AttachmentModel).filter(
            AttachmentModel.file_name.in_(["template.pptx", "001.mp3"])
        ).all()
        self.assertEqual(len(uploaded), 2)
        self.assertTrue(all(item.extracted_text == "" for item in uploaded))
        db.close()


if __name__ == "__main__":
    unittest.main()
